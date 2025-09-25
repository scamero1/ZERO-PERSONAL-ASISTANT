# Zero.py — Typing real, historial visible (selector + botones), scroll auto, dark, sesión 1h, lectura docs
# -----------------------------------------------------------------------------------------------
# Requisitos base:
#   pip install streamlit requests python-dotenv twilio SpeechRecognition av numpy pillow streamlit-webrtc
# Fallbacks (opcionales) y Nombres (OBLIGATORIOS para la seguridad):
#   pip install PyPDF2 python-pptx pandas openpyxl xlrd pytesseract pillow
# -----------------------------------------------------------------------------------------------

import streamlit as st
from streamlit_webrtc import webrtc_streamer, WebRtcMode
import streamlit.components.v1 as components
from PIL import Image, ExifTags
from dotenv import load_dotenv
from twilio.rest import Client
from pathlib import Path
from datetime import datetime
import time, os, json, uuid, requests
import base64 # <-- AÑADIDO: NECESARIO PARA GROQ VISION Y BASE64

# --- Fallback imports protegidos ---
try:
    import PyPDF2
except Exception:
    PyPDF2 = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    import pandas as pd
except Exception:
    pd = None
try:
    import pytesseract
except Exception:
    pytesseract = None
# try:
#     import speech_recognition as sr # Si se implementa transcripción
# except Exception:
#     sr = None

# --- Proyecto (Asume que estas clases ya manejan seguridad (hashing de passwords) ---
# ***********************************************************************************
# IMPORTANTE: Asegúrate de que Login.py use bcrypt o argon2 para hashear contraseñas.
# ***********************************************************************************
from Login import verificar_login, logout as user_logout_logic, registrar_usuario
from database import ZeroDatabase
from file_processor import FileProcessor

# ------------------ Configuración ------------------
load_dotenv()
st.set_page_config(page_title="ZERO - Asistente Virtual", page_icon="favicon.ico", layout="centered")

db = ZeroDatabase()

GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
if not GROQ_API_KEY:
    st.error("Falta GROQ_API_KEY en tu entorno (.env).")
    st.stop()

API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_TEXT_MODEL = os.getenv("GROQ_TEXT_MODEL", "llama-3.1-8b-instant")
GROQ_VISION_MODEL = os.getenv("GROQ_VISION_MODEL", "llama-3.2-11b-vision-preview")
BASE_HEADERS = {"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"}

LOCAL_CHAT_DIR = "local_chats"         # 💾 historial local
SESSION_STORE_DIR = ".auth_sessions"   # 💾 sesión persistente
INACTIVITY_LIMIT_SECS = 60 * 60        # 1 hora
MAX_CONTEXT_MESSAGES = 10              # Límite de mensajes a enviar a la IA para evitar sobrecoste/límite de tokens

# ------------------ Helpers ------------------
def safe_text(x):
    """Mejora de encoding, mantiene la lógica original."""
    if x is None: return ""
    if not isinstance(x, str): x = str(x)
    try: return x.encode("latin1").decode("utf-8")
    except Exception: return x

def _normalize_process_result(res):
    """Se mantiene la lógica de tu extractor original, aunque se recomienda estandarizarlo."""
    content = summary = error = None
    if isinstance(res, dict):
        content = res.get("content") or res.get("text") or res.get("content_extracted")
        summary = res.get("summary") or res.get("analysis") or res.get("brief")
        error = res.get("error")
        return content, summary, error
    if isinstance(res, (list, tuple)):
        if len(res) == 3: content, summary, error = res
        elif len(res) == 2:
            a, b = res
            if isinstance(b, Exception) or (isinstance(b, str) and ("error" in b.lower() or "traceback" in b.lower())):
                content, error = a, b
            else:
                content, summary = a, b
        elif len(res) == 1:
            content = res[0]
        return content, summary, error
    if res is not None: content = str(res)
    return content, summary, error

def clean_session():
    """MEJORA: Centraliza la lógica de limpieza de sesión."""
    ss = st.session_state
    if ss.get("user_id"):
        try: (_sess_file(ss["user_id"])).unlink(missing_ok=True)
        except Exception: pass
    user_logout_logic() # Llama a la lógica de Login.py
    ss.autenticado = False
    ss.usuario = None
    ss.rol = None
    ss.auth_token = None
    ss.user_id = None
    ss.messages = []
    ss.current_chat = str(uuid.uuid4())
    st.rerun()

def _init_state():
    ss = st.session_state
    ss.setdefault("messages", [])
    ss.setdefault("current_chat", str(uuid.uuid4()))
    ss.setdefault("usuario", None)
    ss.setdefault("user_id", None)
    ss.setdefault("rol", None)
    ss.setdefault("autenticado", False)
    ss.setdefault("user_files", [])
    ss.setdefault("user_context", [])
    ss.setdefault("auth_token", None)
    ss.setdefault("last_activity", time.time())
    ss.setdefault("typing_effect", True)    # Typing activado por defecto
    ss.setdefault("selected_local_key", "") # selección en historial local
_init_state()

# ------------------ Estilos (dark + ocultar 3 puntos) ------------------
def load_css():
    """Se mantiene la carga de estilos, excelente base dark."""
    favicon_b64 = ""
    if os.path.exists("favicon.ico"):
        with open("favicon.ico", "rb") as f: favicon_b64 = f.read().hex()
    st.markdown("""
    <style>
    :root {
        color-scheme: dark;
        --bg: #0c0d0f; --card: #15171b; --sidebar: #0f1115;
        --text: #ffffff; --muted: #a0a6b3; --border: #2a2e35;
        --accent: #8B5CF6; --accent2:#7C3AED;
    }
    /* Ocultar menú de 3 puntos y toolbar */
    [data-testid="stMainMenu"], [data-testid="stToolbar"], header { display:none !important; }
    .stApp, [data-testid="stAppViewContainer"], [data-testid="stSidebar"], .block-container {
        background: var(--bg) !important; color: var(--text) !important;
    }
    .chat-box {
        background: var(--card); border:1px solid var(--border); border-radius:12px;
        padding:12px; max-height:72vh; overflow-y:auto;
    }
    /* Mejora de botones para consistencia */
    .stButton>button { background:var(--accent) !important; border:none; color:white; border-radius:10px; font-weight:600; }
    .stButton>button:hover { background:var(--accent2) !important; }
    .stTextInput>div>div>input, .stTextArea textarea, .stSelectbox [data-testid="stNoOptions"], .stSelectbox [data-testid="stSelectboxContainer"] {
        background: var(--card) !important; color: var(--text) !important; border:1px solid var(--border) !important; border-radius:10px !important;
    }
    .sidebar-title {
        font-size:1.1rem; font-weight:700; color:var(--accent); border-bottom:2px solid var(--accent); padding-bottom:6px; margin-bottom:6px;
    }
    .chat-card { background: var(--card); border:1px solid var(--border); border-radius:10px; padding:10px; margin-bottom:8px; }
    </style>
    """, unsafe_allow_html=True)
load_css()

# ------------------ Servicios externos opcionales ------------------
try:
    twilio_client = Client(os.getenv("TWILIO_ACCOUNT_SID"), os.getenv("TWILIO_AUTH_TOKEN"))
except Exception:
    pass

# ------------------ GROQ ------------------
def _system_prompt():
    """Genera el prompt de sistema, incluyendo contexto."""
    base = "Eres un asistente AI llamado Zero. Sé conciso, profesional y útil."
    context_list = st.session_state.get("user_context") or []
    if context_list:
        base += "\n\nContexto personalizado del usuario (archivos relevantes o datos recientes):\n"
        # MEJORA: Solo los 3 contextos más relevantes/recientes para evitar sobrecarga
        for ctx in context_list[-3:]: 
            # Límite a 500 caracteres para el prompt de sistema
            content_snippet = (ctx['context_value'][:500] + "...") if len(ctx['context_value']) > 500 else ctx['context_value']
            base += f"- Clave: {ctx['context_key']}. Contenido:\n{content_snippet}\n"
    return {"role":"system","content":base}

def groq_chat(messages, *, max_tokens=1600, temperature=0.7):
    """Llama a la API de chat de Groq con la configuración de contexto y mensajes."""
    # MEJORA: Truncamiento de mensajes si el historial es muy largo.
    # Se envían el prompt de sistema + los últimos N mensajes (MAX_CONTEXT_MESSAGES).
    history = messages[-(MAX_CONTEXT_MESSAGES - 1):] # Mantiene un mensaje de seguridad
    
    payload = {
        "model": GROQ_TEXT_MODEL, 
        "messages": [_system_prompt()] + history, # Se agrega el system prompt y el historial limitado
        "max_tokens": max_tokens, 
        "temperature": temperature
    }
    
    try:
        r = requests.post(API_URL, headers=BASE_HEADERS, json=payload, timeout=90)
        if r.status_code != 200:
            return f"⚠️ Error {r.status_code}: {r.text}"
        data = r.json()
        return data["choices"][0]["message"]["content"]
    except requests.exceptions.Timeout:
        return "❌ Error: La solicitud a la IA agotó el tiempo de espera (90s). Intenta de nuevo."
    except Exception as e:
        return f"❌ Error de conexión con Groq: {e}"

# ------------------ Sesión persistente (archivo) ------------------
# Lógica mantenida, pero se recuerda que centralizar en BD sería más seguro/escalable.
def _sess_dir(): d = Path(SESSION_STORE_DIR); d.mkdir(parents=True, exist_ok=True); return d
def _sess_file(user_id): return _sess_dir() / f"{user_id}.json"

def persist_login():
    ss = st.session_state
    if not (ss.get("autenticado") and ss.get("user_id")): return
    data = {
        "auth_token": ss.get("auth_token") or str(uuid.uuid4()),
        "user_id": ss["user_id"],
        "usuario": ss["usuario"],
        "rol": ss["rol"],
        "last_activity": time.time()
    }
    ss.auth_token = data["auth_token"]
    _sess_file(ss["user_id"]).write_text(json.dumps(data), encoding="utf-8")

def restore_any_session():
    ss = st.session_state
    if ss.get("autenticado"): return
    try:
        files = list(_sess_dir().glob("*.json"))
        if not files: return
        best, best_ts = None, 0
        for fp in files:
            try:
                data = json.loads(fp.read_text(encoding="utf-8"))
                ts = float(data.get("last_activity", 0))
                if ts > best_ts and (time.time()-ts) <= INACTIVITY_LIMIT_SECS:
                    best, best_ts = data, ts
            except Exception:
                continue
        if best:
            ss.autenticado = True
            ss.user_id = best["user_id"]
            ss.usuario = best.get("usuario")
            ss.rol = best.get("rol")
            ss.auth_token = best.get("auth_token")
            ss.last_activity = time.time()
            # Precarga inicial de archivos y contexto
            ss.user_files = db.get_user_files(ss.user_id)
            ss.user_context = db.get_user_context(ss.user_id)
    except Exception:
        pass

def update_activity():
    ss = st.session_state
    ss.last_activity = time.time()
    if ss.get("user_id"):
        data = {
            "auth_token": ss.get("auth_token") or str(uuid.uuid4()),
            "user_id": ss["user_id"],
            "usuario": ss.get("usuario"),
            "rol": ss.get("rol"),
            "last_activity": ss.last_activity
        }
        _sess_file(ss["user_id"]).write_text(json.dumps(data), encoding="utf-8")

def check_inactivity():
    ss = st.session_state
    if not (ss.get("autenticado") and ss.get("user_id")): return
    try:
        # Esto solo es necesario para verificar si *otro proceso* cerró la sesión
        data = json.loads(_sess_file(ss["user_id"]).read_text(encoding="utf-8"))
        last = float(data.get("last_activity", 0))
        if time.time() - last > INACTIVITY_LIMIT_SECS:
            st.warning("Sesión cerrada por inactividad (1 hora).")
            clean_session() # Llama a la función centralizada
            st.stop()
    except Exception:
        pass

# ------------------ Historial local ------------------
def _user_local_dir(owner):
    d = Path(LOCAL_CHAT_DIR) / str(owner or "anon")
    d.mkdir(parents=True, exist_ok=True)
    return d

def local_save_chat(user_id, chat_id, messages, title="Nuevo chat"):
    d = _user_local_dir(user_id)
    payload = {"chat_id": chat_id, "title": title, "updated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"), "messages": messages}
    (d / f"{chat_id}.json").write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

def local_list_all(user_id):
    items = []
    # MEJORA: Prioriza el historial del usuario loggeado sobre el 'anon'
    owners = [user_id] if user_id else ["anon"] 
    if user_id and user_id != "anon": owners.append("anon")

    for owner in owners:
        d = _user_local_dir(owner)
        for fp in d.glob("*.json"):
            try:
                data = json.loads(fp.read_text(encoding="utf-8"))
                title = data.get("title") or "Nuevo chat"
                updated = data.get("updated_at") or ""
                count = len(data.get("messages") or [])
                # Añadir un prefijo claro
                prefix = "👤" if owner != "anon" else "👻"
                key = f"{owner}||{fp.name}"
                label = f"{prefix} {title} — {updated}  ·  💬 {count}"
                items.append((key, label, str(fp)))
            except Exception:
                continue
    # ordenar por fecha desc
    def _ts(lbl):
        try:
            # Ajustado el split para manejar el nuevo prefijo
            t = lbl.split(" — ")[1].split("  ·")[0].strip()
            return datetime.strptime(t, "%Y-%m-%d %H:%M:%S")
        except Exception:
            return datetime.min
    items.sort(key=lambda x: _ts(x[1]), reverse=True)
    return items

def local_open_chat(path):
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    return data

def local_delete_chat(path):
    Path(path).unlink(missing_ok=True)

# ------------------ Extractores Fallback ------------------
# (Se omiten por espacio, pero se mantienen igual que en el código original)
# ... extract_pdf_text, extract_pptx_text, extract_excel_text, extract_image_text_or_meta, fallback_extract ...
# Se mantiene la lógica original de extracción.

# --- Lógica de Extracción de Archivos (Mantenida) ---

def extract_pdf_text(path):
    if not PyPDF2: return None, "Instala PyPDF2 (pip install PyPDF2)"
    try:
        text = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for p in reader.pages: text.append(p.extract_text() or "")
        return "\n".join(text).strip(), None
    except Exception as e:
        return None, f"PDF error: {e}"

def extract_pptx_text(path):
    if not Presentation: return None, "Instala python-pptx (pip install python-pptx)"
    try:
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape,"has_text_frame") and shape.has_text_frame:
                    parts.append(shape.text)
        return "\n".join(parts).strip(), None
    except Exception as e:
        return None, f"PPTX error: {e}"

def extract_excel_text(path):
    if not pd: return None, "Instala pandas/openpyxl/xlrd"
    try:
        lower = path.lower()
        if lower.endswith(".csv"):
            df = pd.read_csv(path, nrows=200)
        elif lower.endswith(".xlsx"):
            df = pd.read_excel(path, engine="openpyxl")
        elif lower.endswith(".xls"):
            df = pd.read_excel(path)
        else:
            return None, "Formato excel no reconocido"
        # Usar .head() para no sobrecargar la memoria/BD con un archivo gigante
        return df.head(50).to_csv(index=False), None 
    except Exception as e:
        return None, f"Excel error: {e}"

def extract_image_text_or_meta(path):
    try:
        img = Image.open(path)
        txt = ""
        if pytesseract:
            try:
                # Intenta OCR si pytesseract está disponible
                txt = (pytesseract.image_to_string(img) or "").strip()
            except Exception: pass
        
        # Siempre añade metadatos
        meta = [f"Imagen: {os.path.basename(path)}", f"Tamaño: {img.size[0]}x{img.size[1]} px", f"Modo: {img.mode}"]
        try:
            exif = img._getexif()
            if exif:
                tag = {ExifTags.TAGS.get(k, k): v for k,v in exif.items()}
                for k in ("Make","Model","DateTime","Software"):
                    if k in tag: meta.append(f"{k}: {tag[k]}")
        except Exception: pass
        
        final_content = (f"{txt}\n\n---\n[Metadatos]\n" + "\n".join(meta)).strip() if txt else "\n".join(meta).strip()
        
        if not final_content: return None, f"No se pudo extraer texto ni metadatos de {os.path.basename(path)}"
        return final_content, None
    except Exception as e:
        return None, f"Imagen error: {e}"

def fallback_extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf": return extract_pdf_text(path)
    if ext == ".pptx": return extract_pptx_text(path)
    if ext in (".xlsx",".xls",".csv"): return extract_excel_text(path)
    if ext in (".jpg",".jpeg",".png",".gif",".bmp",".webp",".tif",".tiff"): return extract_image_text_or_meta(path)
    return None, f"Sin extractor para {ext}"

# ------------------ Archivos ------------------
def save_uploaded_file(uploaded_file, user_id):
    user_dir = Path("uploads") / str(user_id)
    user_dir.mkdir(parents=True, exist_ok=True)
    dst = user_dir / uploaded_file.name
    dst.write_bytes(uploaded_file.getbuffer())

    processor = FileProcessor()
    
    # Intenta el procesamiento principal
    raw = None
    try:
        # Asumiendo que FileProcessor maneja bien la ruta
        raw = processor.process_file(str(dst), uploaded_file.name)
    except TypeError:
        raw = processor.process_file(str(dst))
    except Exception as e:
        raw = {"error": f"Error del procesador principal: {e}"}

    content, summary, error = _normalize_process_result(raw)
    
    fallback_used = False
    if (not content or content.strip() == "") or error:
        fb_text, fb_err = fallback_extract(str(dst))
        if fb_text:
            content = fb_text if not content else f"{content}\n\n---\n[Fallback]\n{fb_text}"
            summary = summary or f"Contenido extraído de {uploaded_file.name} (vía Fallback)."
            error = None
            fallback_used = True
        elif fb_err and not error:
            error = fb_err
            
    if error: 
        # Si hay error, se intenta borrar el archivo para no dejar basura
        try: dst.unlink(missing_ok=True)
        except Exception: pass
        return None, str(error)

    # Lógica de tipo de archivo (mantenida)
    def _ft(n): return n.split(".")[-1].lower() if "." in n else "txt"
    file_type_getter = getattr(FileProcessor, "get_file_type", _ft)
    file_type = file_type_getter(uploaded_file.name)

    file_id = db.save_file(
        user_id=user_id, filename=uploaded_file.name, file_path=str(dst),
        file_type=file_type, file_size=uploaded_file.size,
        content_extracted=content or "", analysis_summary=summary or ""
    )
    if content:
        db.save_user_context(user_id, f"Archivo: {uploaded_file.name}", content, file_id)
    
    return file_id, "Se utilizó un extractor de reserva." if fallback_used else None

# ------------------ Contexto personalizado ------------------
def get_personalized_context(user_id, query):
    """MEJORA: La lógica de relevancia se mantiene simple para no añadir complejidad con embeddings."""
    try:
        ctxs = db.get_user_context(user_id)
        if not ctxs: return ""
        ql = (query or "").lower().split()
        relev = []
        # Búsqueda por palabras clave (simple)
        for c in ctxs:
            # Solo busca en los últimos 2500 caracteres del contexto para eficiencia
            body = (c.get("context_data") or "").lower()[-2500:] 
            if any(w for w in ql if len(w) > 3 and w in body):
                relev.append((c.get("context_key","Contexto"), c.get("context_data","")))
        
        # Ordenar por el contexto más reciente (los contextos se obtienen en orden de creación)
        relev.reverse() 
        
        if not relev: return ""
        out = "\n\n--- Contexto personalizado de Archivos Relevantes ---\n"
        for k,v in relev[:2]: # Limitar a los 2 más relevantes
            # Limitar el contenido enviado al prompt a 1000 caracteres
            out += f"\n**{k}:**\n{v[:1000]}{'...' if len(v)>1000 else ''}\n"
        out += "--------------------------------------------------------\n"
        return out
    except Exception:
        return ""

# ------------------ Scroll automático ------------------
def _auto_scroll():
    components.html("""
    <script>
    const tryScroll = () => {
        // Selecciona el último chat-box (el principal)
        const els = parent.document.querySelectorAll('.chat-box');
        if (els.length) { 
            const el = els[els.length-1]; 
            // MEJORA: Scroll suave
            el.scrollTo({ top: el.scrollHeight, behavior: 'smooth' });
        }
    };
    // Múltiples intentos para asegurar el scroll después de renderizado/typing
    setTimeout(tryScroll, 50); setTimeout(tryScroll, 200); setTimeout(tryScroll, 600); setTimeout(tryScroll, 1000);
    </script>
    """, height=0)

# ------------------ Chat ------------------
def render_chat(messages):
    st.markdown('<div class="chat-box">', unsafe_allow_html=True)
    for m in messages:
        role = m["role"] if m["role"] in ("user","assistant") else "assistant"
        with st.chat_message(role):
            st.markdown(m["content"])  # soporta bloques ``` de código
            # Mantenido: Tu lógica de timestamp
            st.caption(datetime.now().strftime("%H:%M")) 
    st.markdown('</div>', unsafe_allow_html=True)
    _auto_scroll()

def typewriter_markdown(full_text, delay=0.005):
    """MEJORA: Reducción del delay para mayor fluidez. Se mantiene la lógica por líneas."""
    placeholder = st.empty()
    acc = ""
    # Se aumenta la velocidad
    for line in full_text.split("\n"):
        acc += line + "\n"
        placeholder.markdown(acc)
        # Reducción de delay
        time.sleep(delay) 
    return acc

def chat_page():
    st.title("💬 Chat con Zero")
    render_chat(st.session_state.messages)

    if prompt := st.chat_input("Escribe tu mensaje…"):
        # Añadimos el mensaje del usuario
        st.session_state.messages.append({"role":"user","content":prompt})
        
        # Para el historial que se muestra, usamos el prompt original
        
        # Preparamos contexto personalizado para la llamada a la IA
        enhanced = prompt
        if st.session_state.get("user_id"):
            # Obtiene el contexto relevante
            ctx = get_personalized_context(st.session_state.user_id, prompt) 
            if ctx: enhanced = f"{prompt}\n\n{ctx}"
        
        # El historial de mensajes para la API
        api_messages = [{"role":m["role"],"content":m["content"]} for m in st.session_state.messages[:-1]]
        # El último mensaje es el enhanced prompt con contexto
        api_messages.append({"role":"user","content":enhanced}) 

        # Llamamos al modelo
        with st.spinner("Zero está pensando..."): # MEJORA: Feedback de carga
            raw_reply = groq_chat(api_messages)
        
        reply = safe_text(raw_reply)

        # Mostramos respuesta con efecto typing
        with st.chat_message("assistant"):
            if st.session_state.get("typing_effect", True):
                rendered = typewriter_markdown(reply, delay=0.005) # Delay más bajo
            else:
                st.markdown(reply)
                rendered = reply
            st.caption(datetime.now().strftime("%H:%M"))

        # Guardamos en memoria del chat
        st.session_state.messages.append({"role":"assistant","content":rendered})

        save_chat(); update_activity()
        _auto_scroll()
        st.rerun() # Rerun para limpiar el chat_input

def save_chat():
    """Lógica unificada para guardar en BD y local."""
    try:
        title = "Nuevo chat"
        for m in st.session_state.messages:
            if m["role"] == "user" and m["content"].strip():
                # Toma la primera línea del primer mensaje del usuario
                first_line = m["content"].strip().split('\n')[0] 
                title = (first_line[:50] + "...") if len(first_line)>50 else first_line
                break
        
        # Guardar en BD (si está autenticado)
        if st.session_state.get("user_id"):
            chat_id = db.save_chat(user_id=st.session_state.user_id, chat_id=st.session_state.current_chat, title=title)
            # Borra y reescribe mensajes para mantener el orden correcto en caso de edición/re-guardado
            db.delete_messages_by_chat_id(chat_id) 
            for i, m in enumerate(st.session_state.messages):
                db.save_message(chat_id=chat_id, role=m["role"], content=m["content"], message_order=i)
        
        # Guardar en local (para todos, incluyendo 'anon')
        local_save_chat(st.session_state.get("user_id"), st.session_state.current_chat, st.session_state.messages, title=title)
    except Exception as e:
        print("Error guardando chat:", e)

# ------------------ Sidebar: historial SIMPLE que sí se ve ------------------
def sidebar():
    st.markdown('<div class="sidebar-title">ZERO - Asistente Virtual</div>', unsafe_allow_html=True)
    st.caption(f"Bienvenido, **{st.session_state.get('usuario','Usuario Anónimo')}**")

    # Selector de historial local
    st.subheader("💾 Historial local")
    uid = st.session_state.get("user_id") if st.session_state.get("user_id") is not None else "anon"
    items = local_list_all(uid)
    
    # Crear un índice inicial seguro para el selectbox
    initial_index = 0
    if st.session_state.selected_local_key:
        try:
            # Buscar el índice de la clave seleccionada previamente
            keys = [key for key, _, _ in items]
            initial_index = keys.index(st.session_state.selected_local_key)
        except ValueError:
            st.session_state.selected_local_key = "" # Limpiar si no se encuentra

    if not items:
        st.caption("No hay chats locales aún.")
        selected_key = ""
    else:
        labels = [lbl for _, lbl, _ in items]
        keys = [key for key, _, _ in items]
        paths = {key: path for key, _, path in items}

        selected_label = st.selectbox("Mis chats", labels, index=initial_index, key="local_chat_selector")
        selected_key = keys[labels.index(selected_label)] if labels else ""
        st.session_state.selected_local_key = selected_key # Guarda la clave actual

        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("Abrir", key="open_local", use_container_width=True, disabled=not selected_key):
                data = local_open_chat(paths[selected_key])
                st.session_state.current_chat = data.get("chat_id") or str(uuid.uuid4())
                st.session_state.messages = data.get("messages") or []
                update_activity(); st.rerun()
        with col2:
            if selected_key:
                with open(paths[selected_key], "rb") as fh:
                    st.download_button("Descargar", data=fh.read(),
                                       file_name=os.path.basename(paths[selected_key]),
                                       mime="application/json", use_container_width=True)
            else:
                st.button("Descargar", disabled=True, use_container_width=True)
        with col3:
            if st.button("Eliminar", key="delete_local", use_container_width=True, disabled=not selected_key):
                local_delete_chat(paths[selected_key])
                st.session_state.selected_local_key = "" # Limpiar la selección después de borrar
                update_activity(); st.rerun()

    st.divider()

    # Historial BD (opcional)
    if st.session_state.get("user_id"):
        st.subheader("🗄️ Historial en BD")
        try:
            user_chats = db.get_user_chats(st.session_state.user_id)
            if not user_chats: st.caption("No hay chats en BD.")
            else:
                # Mostrar en orden inverso (más recientes primero)
                titles = [f"{c.get('title','Nuevo chat')} ({c.get('updated_at').split(' ')[0]})" for c in user_chats[::-1]] 
                sel = st.selectbox("Chats en BD", titles, key="db_chat_selector") if titles else None
                if sel:
                    idx = titles.index(sel)
                    chat = user_chats[::-1][idx]
                    if st.button("Abrir chat BD", use_container_width=True):
                        msgs = db.get_chat_messages(chat["chat_id"])
                        st.session_state.current_chat = chat["chat_id"]
                        st.session_state.messages = [{"role":m["role"],"content":m["content"]} for m in msgs]
                        update_activity(); st.rerun()
        except Exception:
            st.caption("Base de datos no disponible o error de conexión.")
    
    st.divider()

    # Navegación según rol (usuario normal solo ve Chat/Subir Archivos)
    st.subheader("🧭 Navegación")
    options = ["Chat Principal", "Subir Archivos"]
    if st.session_state.get("rol") == "admin":
        options += ["Análisis de Imágenes", "Transcripción de Audio", "Registro de Usuarios"]
    choice = st.radio("", options, index=0)

    # Controles de Sesión y Nuevo Chat
    col1, col2 = st.columns(2)
    with col1:
        if st.button("➕ Nuevo Chat", use_container_width=True):
            save_chat() # Guarda el chat actual antes de iniciar uno nuevo
            st.session_state.current_chat = str(uuid.uuid4())
            st.session_state.messages = []
            update_activity(); st.rerun()
    with col2:
        # MEJORA: Llama a la función centralizada de limpieza
        if st.button("🚪 Cerrar sesión", use_container_width=True):
            clean_session() 

    # Switch rápido para activar/desactivar typing
    st.toggle("✍️ Mostrar efecto typing", value=st.session_state.get("typing_effect", True), key="typing_effect")

    return choice

# ------------------ Páginas extra ------------------
def image_page():
    st.title("🖼️ Análisis de Imágenes con Groq Vision")
    up = st.file_uploader("Elige una imagen (JPG, PNG, GIF, etc.)", type=["jpg","jpeg","png","gif","bmp","webp","tif","tiff"])
    
    if up is not None:
        st.image(up, caption="Imagen a analizar", use_column_width=True)
        
        prompt_txt = st.text_input("Instrucción para la IA (ej. 'Describe esta imagen en detalle')", 
                                   value=f"Analiza esta imagen '{up.name}' y proporciona un resumen conciso y profesional, enfocándote en los elementos clave.",
                                   key="vision_prompt")
        
        if st.button("🔍 Analizar Imagen con IA Vision"):
            with st.spinner("Analizando imagen..."):
                try:
                    # 1. MEJORA: Codificación Base64 CORRECTA (binario -> b64 -> string)
                    img_base64 = base64.b64encode(up.getvalue()).decode("utf-8")
                    
                    # 2. MEJORA: Usar el MIME type correcto
                    mime_type = up.type if up.type else "image/jpeg"
                    
                    payload = {
                        "model": GROQ_VISION_MODEL,
                        "messages": [{
                            "role":"user",
                            "content":[
                                {"type":"text","text":prompt_txt},
                                {"type":"image_url","image_url":{"url":f"data:{mime_type};base64,{img_base64}"}}
                            ]
                        }],
                        "max_tokens":1000,"temperature":0.3
                    }
                    
                    r = requests.post(API_URL, headers=BASE_HEADERS, json=payload, timeout=90)
                    
                    if r.status_code == 200:
                        st.success("Análisis completado.")
                        st.markdown(r.json()["choices"][0]["message"]["content"])
                    else:
                        st.error(f"❌ Error {r.status_code}: {r.text}")
                        st.json(r.json())
                
                except requests.exceptions.Timeout:
                    st.error("❌ Tiempo de espera agotado. La IA tardó demasiado en responder.")
                except Exception as e:
                    st.error(f"❌ Ocurrió un error en la solicitud: {e}")

            update_activity()

def audio_page():
    st.title("🎤 Transcripción de Audio")
    st.info("En construcción. Implementa `SpeechRecognition` o un API de transcripción aquí.")
    update_activity()

def register_page():
    st.title("👥 Registro de Usuarios")
    if st.session_state.get("rol") != "admin":
        st.error("❌ Acceso denegado. Solo administradores pueden registrar nuevos usuarios.")
        return
    
    # MEJORA: Añadir la opción de desactivar roles para evitar crear admins accidentalmente
    with st.form("registro_form"):
        st.subheader("Crear Nuevo Usuario")
        user = st.text_input("Nombre de usuario")
        pwd = st.text_input("Contraseña", type="password")
        cpwd = st.text_input("Confirmar contraseña", type="password")
        
        col_rol, col_email = st.columns(2)
        with col_rol:
            rol = st.selectbox("Rol", ["usuario","admin"])
        with col_email:
            email = st.text_input("Email (Opcional)") # Email para futura recuperación/notificaciones
            
        ok = st.form_submit_button("Registrar Usuario")
        
        if ok:
            if not user or not pwd: st.error("El nombre de usuario y la contraseña son obligatorios."); return
            if pwd != cpwd: st.error("Las contraseñas no coinciden."); return
            if len(pwd) < 6: st.error("La contraseña debe tener al menos 6 caracteres."); return
            
            try:
                # Asumiendo que 'registrar_usuario' maneja el hasheo de forma segura
                done = registrar_usuario(user, pwd, rol, email) 
                if done:
                    st.success(f"✅ Usuario '{user}' registrado con el rol de '{rol}'.")
                else:
                    st.error("❌ No se pudo registrar. El usuario ya existe o hubo un error de BD.")
            except Exception as e:
                st.error(f"❌ Error al registrar: {str(e)}")
    
    update_activity()

# ------------------ Subir archivos ------------------
def file_upload_page():
    st.title("📁 Gestión de Archivos y Contexto")
    st.write("Sube documentos (PDF, Excel, PPTX, TXT) o imágenes para que Zero los use como **contexto** en el chat.")

    # Asegurar que el user_id esté cargado
    if not st.session_state.get("usuario"):
        st.error("❌ Error de sesión. Vuelve a iniciar sesión.")
        return

    user_id = db.get_user_id_by_username(st.session_state.usuario)
    if not user_id:
        st.error("❌ No se pudo obtener información del usuario. Intenta cerrar e iniciar sesión.")
        return
    st.session_state.user_id = user_id # Asegura que se actualice

    up = st.file_uploader("Elige un archivo", type=[
        "pdf","docx","doc","txt",
        "xlsx","xls","csv","pptx",
        "jpg","jpeg","png","gif","bmp","webp","tif","tiff"
    ], key="file_uploader")

    if up is not None and st.button("🚀 Procesar Archivo y Añadir Contexto"):
        with st.spinner(f"Procesando '{up.name}'... Esto puede tomar unos segundos."):
            file_id, fb_msg = save_uploaded_file(up, user_id)
            if file_id is None: 
                # El error viene en fb_msg si file_id es None
                st.error(f"❌ {fb_msg}") 
            else:
                success_msg = f"✅ Archivo '{up.name}' procesado y guardado."
                if fb_msg: success_msg += f" (Nota: {fb_msg})" # Si hubo un mensaje de fallback
                st.success(success_msg)
                
                # Recargar el estado de la sesión
                st.session_state.user_files = db.get_user_files(user_id)
                st.session_state.user_context = db.get_user_context(user_id)
                update_activity(); 
                # st.rerun() # Evita el rerun forzado, ya que el st.file_uploader se limpia solo al presionar el botón

    st.subheader("📋 Archivos Subidos")
    # Recargar si es necesario
    if "user_files" not in st.session_state or not st.session_state.user_files:
        st.session_state.user_files = db.get_user_files(user_id)

    if st.session_state.user_files:
        # Mostrar los archivos más recientes primero
        for f in st.session_state.user_files[::-1]: 
            with st.expander(f"📄 **{f['filename']}** ({f['file_type'].upper()})"):
                col1, col2 = st.columns([3,1])
                with col1:
                    st.write(f"**Tamaño:** {f['file_size']/1024:.1f} KB")
                    st.write(f"**Subido:** {f['uploaded_at']}")
                    
                    summary = f.get("analysis_summary")
                    if summary:
                        st.markdown(f"**Resumen:** {summary[:300]}{'...' if len(summary)>300 else ''}")
                    else:
                        st.caption("No hay resumen disponible, solo contenido extraído.")
                        
                with col2:
                    if st.button("🗑️ Eliminar", key=f"del_{f['id']}", use_container_width=True):
                        try:
                            if os.path.exists(f['file_path']): os.remove(f['file_path'])
                        except Exception as e:
                            st.error(f"Error al borrar archivo físico: {e}")
                        db.delete_file(f['id'], user_id)
                        
                        # Recargar y limpiar el contexto
                        st.session_state.user_files = db.get_user_files(user_id)
                        st.session_state.user_context = db.get_user_context(user_id) 
                        update_activity(); st.success("✅ Archivo y contexto eliminado."); st.rerun()
                        
                    if st.button("💬 Usar en Chat", key=f"use_{f['id']}", use_container_width=True):
                        content_extracted = f.get("content_extracted")
                        if content_extracted:
                            # Limitar el snippet a 1500 caracteres para el chat
                            snippet = content_extracted[:1500]
                            prompt = f"📄 **Contenido del Archivo {f['filename']} ({f['file_type'].upper()}):**\n\n{snippet}..."
                            
                            # Añadir el mensaje de usuario al chat actual
                            st.session_state.messages.append({"role":"user","content":prompt})
                            
                            # Guardar y cambiar a la página principal
                            save_chat(); 
                            st.session_state.selected_page = "Chat Principal" # Establecer la navegación
                            update_activity(); 
                            st.success("Contenido agregado al chat. ¡Vuelve al chat principal!")
                            st.rerun()
                        else:
                            st.warning("El archivo no tiene contenido de texto extraído.")
    else:
        st.info("No hay archivos subidos aún.")

# ------------------ Main ------------------
def main():
    # 1) Expiración e intento de restaurar sesión
    check_inactivity()
    restore_any_session()

    # 2) Si no autenticado, login
    if not st.session_state.get("autenticado", False):
        verificar_login()
        if st.session_state.get("autenticado") and st.session_state.get("user_id"):
            persist_login(); update_activity()
        else:
            return # Detener la ejecución si no está autenticado
    else:
        update_activity()

    # 3) Precarga y mantenimiento de estado (se movió en parte a restore_any_session)
    if st.session_state.get("user_id") and not st.session_state.get("user_files"):
        st.session_state.user_files = db.get_user_files(st.session_state.user_id)
        st.session_state.user_context = db.get_user_context(st.session_state.user_id)

    # 4) Sidebar y routing
    with st.sidebar:
        choice = sidebar()
    
    # MEJORA: Almacenar la elección de navegación para usarla en reruns (ej. al abrir chat)
    st.session_state.selected_page = choice 

    if st.session_state.selected_page == "Chat Principal":
        chat_page()
    elif st.session_state.selected_page == "Subir Archivos":
        file_upload_page()
    elif st.session_state.selected_page == "Análisis de Imágenes":
        if st.session_state.get("rol") == "admin": image_page()
        else: st.error("Acceso denegado.")
    elif st.session_state.selected_page == "Transcripción de Audio":
        if st.session_state.get("rol") == "admin": audio_page()
        else: st.error("Acceso denegado.")
    elif st.session_state.selected_page == "Registro de Usuarios":
        if st.session_state.get("rol") == "admin": register_page()
        else: st.error("Acceso denegado.")

if __name__ == "__main__":
    main()